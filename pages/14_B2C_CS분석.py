import streamlit as st
import pandas as pd
import gspread
from oauth2client.service_account import ServiceAccountCredentials
import plotly.express as px
import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

# ==========================================
# CSS 설정
# ==========================================
import src.style_utils as style
st.set_page_config(page_title="독서화랑 대시보드", layout="wide")
style.apply_common_style()

# ==========================================
# [설정] 페이지 설정
# ==========================================
st.set_page_config(page_title="일반 CS 통합 대시보드", page_icon="📞", layout="wide")
SHEET_URL = "https://docs.google.com/spreadsheets/d/1MQVn2jcKiHagQqUyyHR3ew9BLhD520Cv3UTwVMo5_6g/edit?usp=sharing"

@st.cache_data(ttl=60)
def load_all_data():
    try:
        scope = ['https://spreadsheets.google.com/feeds', 'https://www.googleapis.com/auth/drive']
        creds_dict = dict(st.secrets["gcp_service_account"])
        creds = ServiceAccountCredentials.from_json_keyfile_dict(creds_dict, scope)
        client = gspread.authorize(creds)
        sh = client.open_by_url(SHEET_URL)
        
        def fetch(sheet_name, prefix):
            try:
                ws = sh.worksheet(sheet_name)
                data = ws.get_all_values()
                if len(data) < 5: return pd.DataFrame()
                df = pd.DataFrame(data[5:], columns=[c.strip() for c in data[4]])
                df = df[df['일시'].str.strip() != ''].dropna(subset=['일시'])
                for c in ['일시', '처리일']:
                    if c in df.columns:
                        df[c] = pd.to_datetime(df[c].str.replace('.', '-'), errors='coerce')
                if '처리일' in df.columns:
                    df['체류시간'] = (df['처리일'] - df['일시']).dt.total_seconds() / (86400)
                
                day_map = {0:'월', 1:'화', 2:'수', 3:'목', 4:'금', 5:'토', 6:'일'}
                df['요일'] = df['일시'].dt.dayofweek.map(day_map)
                
                df['데이터소스'] = prefix
                df['고유ID'] = prefix[0] + "_" + df.index.astype(str)
                return df
            except: return pd.DataFrame()

        df_merged = pd.concat([fetch("CS 접수기록(관리부)", "관리부"), fetch("CS 접수기록(선생님)", "선생님")], ignore_index=True)
        return df_merged
    except Exception as e:
        st.error(f"데이터 로드 오류: {e}")
        return pd.DataFrame()

df_raw = load_all_data()
def create_cs_ppt_card(row, name_col, branch_col, q_text, a_text, dept_info, collab_dept, status_bg):

    # 1. PPT 객체 생성 (16:9 비율)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # 빈 슬라이드

    # 2. 메인 카드 배경
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(0.5), Inches(0.5), Inches(12.333), Inches(6.5)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.line.color.rgb = RGBColor(224, 224, 224)

    # 3. 태그 생성 헬퍼 함수
    def add_tag(text, color_hex, left_inch, top_inch, width_inch=1.1):
        r, g, b = tuple(int(color_hex.lstrip('#')[i:i+2], 16) for i in (0, 2, 4))
        tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_inch), Inches(top_inch), Inches(width_inch), Inches(0.35))
        tag.fill.solid()
        tag.fill.fore_color.rgb = RGBColor(r, g, b)
        tag.line.fill.background()
        tf = tag.text_frame
        tf.text = str(text)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        return left_inch + width_inch + 0.1

    # 상단 태그 배치
    l = 0.9
    l = add_tag(row.get('D-1', '-'), "#8E44AD", l, 0.8)
    l = add_tag(row.get('D-2', '-'), "#00897B", l, 0.8, 1.3)
    l = add_tag(row.get('D-3', '-'), "#E67E22", l, 0.8, 1.3)
    l = add_tag(f"{row[name_col]} 학생", "#1976D2", l, 0.8, 1.3)

    # 카테고리 제목 및 메타 정보
    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.25), Inches(5), Inches(0.5))
    title_box.text_frame.text = str(row.get('카테고리', '미분류'))
    title_box.text_frame.paragraphs[0].font.size = Pt(20)
    title_box.text_frame.paragraphs[0].font.bold = True

    meta_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(8), Inches(0.3))
    meta_box.text_frame.text = f"📅 접수: {row['일시'].strftime('%Y-%m-%d %H:%M')} | 🏢 소속: {dept_info} | 🎓 학년: {row.get('학년','-')}"
    meta_box.text_frame.paragraphs[0].font.size = Pt(12)
    meta_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(136, 136, 136)

    # 4. 본문 영역 (좌/우 박스)
    def add_content_box(title, content, bg_hex, line_hex, left_inch):
        bg_r, bg_g, bg_b = tuple(int(bg_hex.lstrip('#')[i:i+2], 16) for i in (0, 2, 4))
        ln_r, ln_g, ln_b = tuple(int(line_hex.lstrip('#')[i:i+2], 16) for i in (0, 2, 4))
        
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left_inch), Inches(2.4), Inches(5.5), Inches(3.2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(bg_r, bg_g, bg_b)
        box.line.fill.background()

        top_line = slide.shapes.add_connector(1, Inches(left_inch), Inches(2.4), Inches(left_inch + 5.5), Inches(2.4))
        top_line.line.color.rgb = RGBColor(ln_r, ln_g, ln_b)
        top_line.line.width = Pt(4)

        t_box = slide.shapes.add_textbox(Inches(left_inch + 0.1), Inches(2.5), Inches(5), Inches(0.4))
        t_box.text_frame.text = title
        t_box.text_frame.paragraphs[0].font.bold = True
        t_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(ln_r, ln_g, ln_b)

        c_box = slide.shapes.add_textbox(Inches(left_inch + 0.1), Inches(2.9), Inches(5.3), Inches(2.6))
        c_box.text_frame.word_wrap = True
        c_box.text_frame.text = str(content)
        c_box.text_frame.paragraphs[0].font.size = Pt(12)

    add_content_box("🗣️ 접수/문의 내용", q_text, "#FDF7F2", "#FF8C00", 0.9)
    add_content_box("🛠️ 실제 처리/답변 내용", a_text, "#F2F9F2", "#2ECC71", 6.7)

    # 5. 하단 푸터
    footer_text = f"최종 처리 결과: {row.get('처리카테고리', '-')}   |   🤝 협업 부서: {collab_dept}   |   ID: {row['고유ID']}   |   처리 소요: {round(row.get('체류시간', 0), 1)}일"
    footer_box = slide.shapes.add_textbox(Inches(0.9), Inches(6.0), Inches(11.5), Inches(0.4))
    footer_box.text_frame.text = footer_text
    footer_box.text_frame.paragraphs[0].font.size = Pt(11)
    footer_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

    # 파일 저장
    file_path = "temp_report.pptx"
    prs.save(file_path)
    return file_path

# ==========================================
# [UI] 사이드바 및 필터
# ==========================================
with st.sidebar:
    st.title("🗂️ 조회 모드")
    target_mode = st.radio("데이터 범위", ["전체(통합)", "관리부", "선생님"])
    st.divider()
    start_date = st.date_input("시작일", datetime.date(2025, 12, 3))
    end_date = st.date_input("종료일", datetime.date.today())

start_dt, end_dt = pd.to_datetime(start_date), pd.to_datetime(end_date)
mask = (df_raw['일시'] >= start_dt) & (df_raw['일시'] <= end_dt)
if target_mode != "전체(통합)": mask &= (df_raw['데이터소스'] == target_mode)
df = df_raw.loc[mask].copy()

if df.empty:
    st.warning("데이터가 없습니다.")
    st.stop()

# 메인 타이틀
st.title(f"📞 일반 CS 통합 분석 [{target_mode}]")

# 탭 생성
tab1, tab2,tab3 = st.tabs(["📊 종합 분석 리포트", "🔎 개별 상세 조회","🔎 유형별 상세 조회"])

# ==========================================
# [탭 1] 종합 분석 리포트
# ==========================================
with tab1:
    # KPI 지표
    k1, k2, k3 = st.columns(3)
    k1.metric("총 접수", f"{len(df)}건")
    k2.metric("평균 처리 시간", f"{df['체류시간'].mean():.1f}일")
    k3.metric("최다 발생 이슈", df['카테고리'].value_counts().idxmax() if '카테고리' in df.columns else "-")

    st.divider()

    # 1. 히트맵 분석
    st.subheader("1. 🗺️ 접수-처리 집중도 분석 (Heatmap)")
    if '카테고리' in df.columns and '처리카테고리' in df.columns:
        ct = pd.crosstab(df['카테고리'], df['처리카테고리'])
        fig_heat = px.imshow(ct, text_auto=True, aspect="auto", color_continuous_scale="Reds")
        st.plotly_chart(fig_heat, use_container_width=True, key="chart_1")

    st.divider()

    # 2. 분포 분석
    st.subheader("2. 항목별 상세 분포")
    col_a, col_b = st.columns(2)

    with col_a:
        st.markdown("##### 🏢 부서별 이슈 관여도")
        if '협업 부서' in df.columns:
            # 공백 제거 및 필터링된 df 사용 확인
            # value_counts()는 현재 df(필터링된 상태)의 분포를 그대로 반영합니다.
            dept_cnt = df[df['협업 부서'].str.strip().fillna('') != '']['협업 부서'].value_counts().reset_index()
            dept_cnt.columns = ['협업 부서', '건수'] # 컬럼명 명확히 지정
            
            fig_dept = px.bar(dept_cnt, x='건수', y='협업 부서', orientation='h', 
                            color_discrete_sequence=['#FF8C00'],
                            category_orders={"협업 부서": dept_cnt['협업 부서'].tolist()}) # 빈도순 정렬 유지
            st.plotly_chart(fig_dept, use_container_width=True, key="chart_2")

    with col_b:
        st.markdown("##### 📅 요일별 접수량")
        if '요일' in df.columns:
            day_order = ['월', '화', '수', '목', '금', '토', '일']
            # 필터링된 df에서 요일 빈도 계산
            day_cnt = df['요일'].value_counts().reindex(day_order).fillna(0).reset_index()
            day_cnt.columns = ['요일', '건수']
            
            fig_day = px.bar(day_cnt, x='요일', y='건수', color_discrete_sequence=['#A9A9A9'])
            st.plotly_chart(fig_day, use_container_width=True, key="chart_4")

    # 3. 리스크 진단
    st.subheader("3. 🚨 서비스 안정성 진단")
    def classify_risk(val):
        val = str(val).strip()
        if val in ['시스템오류', '회원연동문제']: return '⛔ 심각 오류 (시스템/연동)'
        elif val == '컨텐츠오류': return '📉 컨텐츠오류'
        else: return '⚠️ 일반문의/기타'

    target_col = '처리카테고리' if '처리카테고리' in df.columns else '카테고리'
    df['리스크'] = df[target_col].apply(classify_risk)
    risk_cnt = df['리스크'].value_counts().reset_index(name='건수')
    fig_risk = px.pie(risk_cnt, values='건수', names='리스크', hole=0.5, 
                     color='리스크',
                     color_discrete_map={'⛔ 심각 오류 (시스템/연동)': '#FF4B4B', '📉 컨텐츠오류': '#FF8C00', '⚠️ 일반문의/기타': '#E0E0E0'})
    st.plotly_chart(fig_risk, use_container_width=True, key="chart_6")

    # ---------------------------------------------------------
    # 추가 요청: 📈 날짜별 접수량 추이 (CS 감소 추세 확인용)
    # ---------------------------------------------------------
    st.divider()
    st.subheader("3. 📈 일자별 CS 접수량 추이")
    if '일시' in df.columns:
        # 날짜별로 건수 집계
        daily_trend = df.groupby(df['일시'].dt.date).size().reset_index(name='접수건수')
        daily_trend.columns = ['날짜', '접수건수']
        
        # 라인 그래프 생성
        fig_trend = px.line(daily_trend, x='날짜', y='접수건수', 
                            markers=True,
                            line_shape="linear",
                            color_discrete_sequence=['#EF553B']) # 눈에 띄는 색상
        
        # 레이아웃 미세 조정 (날짜 포맷 등)
        fig_trend.update_layout(hovermode="x unified")
        st.plotly_chart(fig_trend, use_container_width=True, key="chart_trend")
       

# ==========================================
# [탭 2] 개별 상세 조회 (기존 탭 유지)
# ==========================================
with tab2:
    st.subheader("🔍 개별 건 상세 정보 열람")

    # 1. 필수 컬럼명 매칭 
    name_col = '이름' if '이름' in df.columns else ('학생명' if '학생명' in df.columns else '이름')
    branch_col = '지점' if '지점' in df.columns else ('소속지점' if '소속지점' in df.columns else '지점')
    collab_col = '협업 부서' if '협업 부서' in df.columns else ('협업부서' if '협업부서' in df.columns else '협업 부서')
    q_col = '문의 내용' if '문의 내용' in df.columns else '문의내용'
    a_col = '처리 내용' if '처리 내용' in df.columns else '처리내용'

    # 2. 조회 방식 선택
    search_method = st.radio("조회 방식 선택", ["카테고리 필터로 찾기", "학생 이름으로 검색"], horizontal=True)

    selected_row = None

    if search_method == "학생 이름으로 검색":
        search_name = st.text_input("👤 조회할 학생 이름을 입력하세요", placeholder="이름 입력 후 엔터")
        if search_name:
            name_history = df[df[name_col].str.strip() == search_name.strip()].copy()
            if not name_history.empty:
                st.success(f"✅ '{search_name}' 학생의 데이터 총 **{len(name_history)}**건을 찾았습니다.")
                name_history['display_label'] = (
                    name_history['일시'].dt.strftime('%Y-%m-%d') + " | " + 
                    name_history[branch_col].astype(str) + " | " +
                    name_history['카테고리']
                )
                selected_item = st.selectbox("📑 열람할 상담 건을 선택하세요", name_history['display_label'].tolist())
                selected_row = name_history[name_history['display_label'] == selected_item].iloc[0]
            else:
                st.error(f"❌ '{search_name}' 학생의 데이터가 조회 기간 내에 없습니다.")
    else:
        # 필터 UI
        f1, f2 = st.columns(2)
        with f1: 
            sel_in = st.selectbox("📥 접수 유형 필터", ["전체"] + sorted(df['카테고리'].astype(str).unique().tolist()))
        with f2: 
            sel_out = st.selectbox("📤 처리 결과 필터", ["전체"] + sorted(df['처리카테고리'].astype(str).unique().tolist()) if '처리카테고리' in df.columns else ["전체"])

        # 필터링 로직
        df_filtered = df.copy()
        if sel_in != "전체":
            df_filtered = df_filtered[df_filtered['카테고리'] == sel_in]
        if sel_out != "전체" and '처리카테고리' in df_filtered.columns:
            df_filtered = df_filtered[df_filtered['처리카테고리'] == sel_out]

        st.info(f"✅ 현재 조건으로 검색된 데이터: 총 **{len(df_filtered)}**건")

        if not df_filtered.empty:
            st.dataframe(df_filtered[['고유ID', '일시', branch_col, name_col, '카테고리', '처리 상태']].sort_values('일시', ascending=False), 
                         use_container_width=True, hide_index=True, height=200)
            
            st.markdown("---")
            selected_id = st.selectbox("📑 상세 정보를 열람할 '고유ID'를 선택하세요", ["선택하세요"] + df_filtered['고유ID'].tolist())
            if selected_id != "선택하세요":
                selected_row = df_filtered[df_filtered['고유ID'] == selected_id].iloc[0]

    # 3. 상세 카드 출력 영역 
    if selected_row is not None:
        row = selected_row
        status_bg = "#E1F5FE" if row.get('처리 상태') == "처리완료" else "#FFF9C4"
        
        q_text = str(row.get(q_col, '내용 없음')).replace("#", "＃").replace("[", "［").replace("]", "］")
        a_text = str(row.get(a_col, '기록 없음')).replace("#", "＃").replace("[", "［").replace("]", "］")
        
        dept_info = f"{row.get('데이터소스', '-')} ({row.get(branch_col, '-')})"
        collab_val = str(row.get(collab_col, '')).strip()
        collab_dept = collab_val if collab_val not in ['', 'nan', 'None'] else '없음'

        card_html = f"""
        <div style="background-color: white; border: 1px solid #E0E0E0; border-radius: 12px; padding: 25px; box-shadow: 0 4px 10px rgba(0,0,0,0.05); margin-top: 15px; font-family: sans-serif; text-align: left;">
            <div style="display: flex; justify-content: space-between; align-items: start; margin-bottom: 20px; border-bottom: 2px solid #F5F5F5; padding-bottom: 15px;">
                <div>
                    <span style="background-color: #1976D2; color: white; padding: 3px 12px; border-radius: 4px; font-weight: bold; font-size: 13px;">{row[name_col]} 학생</span>
                    <h3 style="margin: 10px 0 5px 0; color: #333; font-size: 18px; font-weight: bold; border:none;">{row.get('카테고리', '미분류')}</h3>
                    <p style="margin: 0; color: #888; font-size: 13px;">📅 접수: {row['일시'].strftime('%Y-%m-%d %H:%M')} | 🏢 소속: <b>{dept_info}</b> | 🎓 학년: {row.get('학년','-')}</p>
                </div>
                <div style="background-color: {status_bg}; padding: 8px 25px; border-radius: 50px; font-weight: bold; color: #444; border: 1px solid #DDD; font-size: 14px;">{row.get('처리 상태', '상태없음')}</div>
            </div>
            <div style="display: grid; grid-template-columns: 1fr 1fr; gap: 20px;">
                <div style="background-color: #FDF7F2; padding: 18px; border-radius: 8px; border-top: 4px solid #FF8C00;">
                    <strong style="color: #D35400; font-size: 14px;">🗣️ 접수/문의 내용</strong>
                    <div style="margin-top: 12px; line-height: 1.6; font-size: 14px !important; color: #444; min-height: 120px; white-space: pre-wrap; font-weight: 400;">{q_text}</div>
                </div>
                <div style="background-color: #F2F9F2; padding: 18px; border-radius: 8px; border-top: 4px solid #2ECC71;">
                    <strong style="color: #27AE60; font-size: 14px;">🛠️ 실제 처리/답변 내용</strong>
                    <div style="margin-top: 12px; line-height: 1.6; font-size: 14px !important; color: #444; min-height: 120px; white-space: pre-wrap; font-weight: 400;">{a_text}</div>
                </div>
            </div>
            <div style="margin-top: 20px; padding-top: 15px; border-top: 1px dashed #DDD; display: flex; justify-content: space-between; color: #666; font-size: 13px;">
                <div style="display: flex; gap: 20px;">
                    <div><b>최종 처리 유형:</b> <span style="color: #1976D2;">{row.get('처리카테고리', '미분류')}</span></div>
                    <div><b>🤝 협업 부서:</b> <span style="color: #E65100; font-weight: bold;">{collab_dept}</span></div>
                    <div><b>🏷️ 분류 유형:</b> <span style="color: #8E44AD; font-weight: bold;">{row.get('유형', '-')}</span></div>
                </div>
                <div><b>ID:</b> {row['고유ID']} | <b>처리 소요:</b> {round(row.get('체류시간', 0), 1)}일</div>
            </div>
        </div>
        """
        st.markdown(card_html, unsafe_allow_html=True)

with tab3:
    st.subheader("🔍 유형별 건 상세 정보 열람")
    
    if 'D-1' in df.columns and 'D-2' in df.columns:
        if 'D-3' not in df.columns:
            df['D-3'] = "-"
        df['D-3'] = df['D-3'].fillna("-").replace("", "-")

        # ==========================================
        # [추가] D-1 순서 강제 지정 (시스템 -> 정책 -> 건의사항)
        # ==========================================
        d1_target_order = ["시스템", "정책", "건의사항"]
        # 혹시 위 3개 외에 다른 분류가 데이터에 있다면 뒤에 자연스럽게 붙도록 처리
        current_d1 = df['D-1'].dropna().unique().tolist()
        d1_order = d1_target_order + [x for x in current_d1 if x not in d1_target_order]
        
        # DataFrame을 Categorical 타입으로 변경하여 정렬 순서를 커스텀 오더로 고정
        df['D-1'] = pd.Categorical(df['D-1'], categories=d1_order, ordered=True)

        # ==========================================
        # 1. 표와 그래프 (전체 너비 단일 구성, 표 3뎁스)
        # ==========================================
        
        # --- 표 영역 (D-1 순서 적용) ---
        type_cat_cnt_3d = df.groupby(['D-1', 'D-2', 'D-3'], observed=True).size().reset_index(name='건수')
        total_cnt = len(df)
        
        if total_cnt > 0:
            type_cat_cnt_3d['비중(%)'] = (type_cat_cnt_3d['건수'] / total_cnt * 100).round(1).astype(str) + "%"
        else:
            type_cat_cnt_3d['비중(%)'] = "0%"
        
        # Categorical이 적용되었으므로 D-1은 시스템->정책->건의사항 순으로 정렬됨
        type_cat_cnt_3d = type_cat_cnt_3d.sort_values(by=['D-1', 'D-2', '건수'], ascending=[True, True, False])
        
        st.markdown("##### 📋 분류별 발생 비중 (D-1 > D-2 > D-3)")
        st.dataframe(type_cat_cnt_3d.set_index(['D-1', 'D-2', 'D-3']), use_container_width=True)
        
        st.markdown("<br>", unsafe_allow_html=True)

        # --- 그래프 영역 (D-1 순서 및 라벨 적용) ---
        st.markdown("##### 📊 대분류별 중분류 누적 막대 그래프")
        
        type_cat_cnt_2d = df.groupby(['D-1', 'D-2'], observed=True).size().reset_index(name='건수')
        
        type_cat_cnt_2d['라벨'] = type_cat_cnt_2d.apply(
            lambda x: f"<b>{x['D-2']}</b><br>{x['건수']}건" if x['건수'] > 0 else "", axis=1
        )
        
        fig_stacked = px.bar(
            type_cat_cnt_2d, 
            x='D-1', 
            y='건수', 
            color='D-2', 
            text='라벨',
            barmode='stack',
            category_orders={"D-1": d1_order}, # 그래프 X축 순서 강제 적용
            color_discrete_sequence=px.colors.qualitative.Set2
        )
        
        fig_stacked.update_traces(
            textposition='inside', 
            insidetextanchor='middle',
            textfont=dict(color='white', size=14)
        )
        
        total_by_class = type_cat_cnt_2d.groupby('D-1', observed=True)['건수'].sum().reset_index()
        total_by_class = total_by_class[total_by_class['건수'] > 0] # 건수가 0인 것은 합계 라벨 제외
        
        if not total_by_class.empty:
            fig_stacked.add_scatter(
                x=total_by_class['D-1'], 
                y=total_by_class['건수'], 
                text='<b>총 ' + total_by_class['건수'].astype(str) + '건</b>',
                mode='text', 
                textposition='top center', 
                textfont=dict(size=16, color='#212529'),
                showlegend=False,
                hoverinfo='skip'
            )
        
        max_y = total_by_class['건수'].max() * 1.15 if not total_by_class.empty else 10
        
        fig_stacked.update_layout(
            xaxis_title="D-1 (대분류)", 
            yaxis_title="접수 건수",
            yaxis=dict(range=[0, max_y]),
            legend_title_text="D-2 (중분류)",
            margin=dict(l=20, r=20, t=30, b=20)
        )
        st.plotly_chart(fig_stacked, use_container_width=True, key="tab3_stacked_chart")
            
        st.divider()

        # ==========================================
        # 2. 유형별 상세 검색 (필터 순서 적용)
        # ==========================================
        st.subheader("🔍 유형별 유형 상세 기록 검색")
        
        df_search_base = df_raw.copy()
        if target_mode != "전체(통합)":
            df_search_base = df_search_base[df_search_base['데이터소스'] == target_mode]

        if 'D-3' not in df_search_base.columns:
            df_search_base['D-3'] = "-"
        df_search_base['D-3'] = df_search_base['D-3'].fillna("-").replace("", "-")

        # 검색용 데이터베이스에도 순서 고정 적용
        df_search_base['D-1'] = pd.Categorical(df_search_base['D-1'], categories=d1_order, ordered=True)

        name_col = '이름' if '이름' in df_search_base.columns else ('학생명' if '학생명' in df_search_base.columns else '이름')
        branch_col = '지점' if '지점' in df_search_base.columns else ('소속지점' if '소속지점' in df_search_base.columns else '지점')
        collab_col = '협업 부서' if '협업 부서' in df_search_base.columns else ('협업부서' if '협업부서' in df_search_base.columns else '협업 부서')
        q_col = '문의 내용' if '문의 내용' in df_search_base.columns else '문의내용'
        a_col = '처리 내용' if '처리 내용' in df_search_base.columns else '처리내용'

        col_f1, col_f2, col_f3 = st.columns(3)
        
        with col_f1:
            # 가나다순(sorted) 대신 지정된 순서대로 필터 목록 생성
            d1_list = [x for x in d1_order if x in df_search_base['D-1'].dropna().unique()]
            sel_d1 = st.selectbox("📌 D-1 (대분류)", ["전체"] + d1_list)
        
        with col_f2:
            if sel_d1 != "전체":
                d2_list = df_search_base[df_search_base['D-1'] == sel_d1]['D-2'].dropna().unique().tolist()
            else:
                d2_list = df_search_base['D-2'].dropna().unique().tolist()
            sel_d2 = st.selectbox("🔖 D-2 (중분류)", ["전체"] + sorted(d2_list))
            
        with col_f3:
            mask_d3 = pd.Series(True, index=df_search_base.index)
            if sel_d1 != "전체": mask_d3 &= (df_search_base['D-1'] == sel_d1)
            if sel_d2 != "전체": mask_d3 &= (df_search_base['D-2'] == sel_d2)
            
            d3_list = df_search_base[mask_d3]['D-3'].dropna().unique().tolist()
            sel_d3 = st.selectbox("📝 D-3 (소분류)", ["전체"] + sorted(d3_list))

        df_type_filtered = df_search_base.copy()
        if sel_d1 != "전체": df_type_filtered = df_type_filtered[df_type_filtered['D-1'] == sel_d1]
        if sel_d2 != "전체": df_type_filtered = df_type_filtered[df_type_filtered['D-2'] == sel_d2]
        if sel_d3 != "전체": df_type_filtered = df_type_filtered[df_type_filtered['D-3'] == sel_d3]

        st.info(f"✅ 현재 선택된 유형의 데이터 (전체 기간): 총 **{len(df_type_filtered)}**건")

        if not df_type_filtered.empty:
            st.dataframe(
                df_type_filtered[['고유ID', '일시', branch_col, name_col, 'D-1', 'D-2', 'D-3', '처리 상태']].sort_values('일시', ascending=False), 
                use_container_width=True, 
                hide_index=True, 
                height=200
            )
            
            st.markdown("---")
            selected_type_id = st.selectbox("📑 상세 기록을 열람할 '고유ID'를 선택하세요", ["선택하세요"] + df_type_filtered['고유ID'].tolist())
            
            if selected_type_id != "선택하세요":
                row = df_type_filtered[df_type_filtered['고유ID'] == selected_type_id].iloc[0]
                status_bg = "#E1F5FE" if row.get('처리 상태') == "처리완료" else "#FFF9C4"
                
                q_text = str(row.get(q_col, '내용 없음')).replace("#", "＃").replace("[", "［").replace("]", "］")
                a_text = str(row.get(a_col, '기록 없음')).replace("#", "＃").replace("[", "［").replace("]", "］")
                
                dept_info = f"{row.get('데이터소스', '-')} ({row.get(branch_col, '-')})"
                collab_val = str(row.get(collab_col, '')).strip()
                collab_dept = collab_val if collab_val not in ['', 'nan', 'None'] else '없음'

                card_html = f"""
                <div style="background-color: white; border: 1px solid #E0E0E0; border-radius: 12px; padding: 25px; box-shadow: 0 4px 10px rgba(0,0,0,0.05); margin-top: 15px; font-family: sans-serif; text-align: left;">
                    <div style="display: flex; justify-content: space-between; align-items: start; margin-bottom: 20px; border-bottom: 2px solid #F5F5F5; padding-bottom: 15px;">
                        <div>
                            <span style="background-color: #8E44AD; color: white; padding: 3px 12px; border-radius: 4px; font-weight: bold; font-size: 13px;">{row.get('D-1', '-')}</span>
                            <span style="background-color: #00897B; color: white; padding: 3px 12px; border-radius: 4px; font-weight: bold; font-size: 13px; margin-left: 5px;">{row.get('D-2', '-')}</span>
                            <span style="background-color: #E67E22; color: white; padding: 3px 12px; border-radius: 4px; font-weight: bold; font-size: 13px; margin-left: 5px;">{row.get('D-3', '-')}</span>
                            <span style="background-color: #1976D2; color: white; padding: 3px 12px; border-radius: 4px; font-weight: bold; font-size: 13px; margin-left: 5px;">{row[name_col]} 학생</span>
                            <h3 style="margin: 10px 0 5px 0; color: #333; font-size: 18px; font-weight: bold; border:none;">{row.get('카테고리', '미분류')}</h3>
                            <p style="margin: 0; color: #888; font-size: 13px;">📅 접수: {row['일시'].strftime('%Y-%m-%d %H:%M')} | 🏢 소속: <b>{dept_info}</b> | 🎓 학년: {row.get('학년','-')}</p>
                        </div>
                        <div style="background-color: {status_bg}; padding: 8px 25px; border-radius: 50px; font-weight: bold; color: #444; border: 1px solid #DDD; font-size: 14px;">{row.get('처리 상태', '상태없음')}</div>
                    </div>
                    <div style="display: grid; grid-template-columns: 1fr 1fr; gap: 20px;">
                        <div style="background-color: #FDF7F2; padding: 18px; border-radius: 8px; border-top: 4px solid #FF8C00;">
                            <strong style="color: #D35400; font-size: 14px;">🗣️ 접수/문의 내용</strong>
                            <div style="margin-top: 12px; line-height: 1.6; font-size: 14px !important; color: #444; min-height: 120px; white-space: pre-wrap; font-weight: 400;">{q_text}</div>
                        </div>
                        <div style="background-color: #F2F9F2; padding: 18px; border-radius: 8px; border-top: 4px solid #2ECC71;">
                            <strong style="color: #27AE60; font-size: 14px;">🛠️ 실제 처리/답변 내용</strong>
                            <div style="margin-top: 12px; line-height: 1.6; font-size: 14px !important; color: #444; min-height: 120px; white-space: pre-wrap; font-weight: 400;">{a_text}</div>
                        </div>
                    </div>
                    <div style="margin-top: 20px; padding-top: 15px; border-top: 1px dashed #DDD; display: flex; justify-content: space-between; color: #666; font-size: 13px;">
                        <div style="display: flex; gap: 20px;">
                            <div><b>최종 처리 결과:</b> <span style="color: #1976D2;">{row.get('처리카테고리', '미분류')}</span></div>
                            <div><b>🤝 협업 부서:</b> <span style="color: #E65100; font-weight: bold;">{collab_dept}</span></div>
                        </div>
                        <div><b>ID:</b> {row['고유ID']} | <b>처리 소요:</b> {round(row.get('체류시간', 0), 1)}일</div>
                    </div>
                </div>
                """
                st.markdown(card_html, unsafe_allow_html=True)
    else:
        st.warning("⚠️ 구글 시트에 'D-1', 'D-2' 컬럼이 존재하지 않습니다. 시트의 컬럼명을 확인해주세요.")

    st.divider() # 시각적 구분선
        
    # 기존 버튼 코드를 아래와 같이 수정 (괄호 안에 변수들을 순서대로 넣어줍니다)
    if st.button("📄 이 상세 카드를 PPT로 내보내기"):
        try:
            # 함수 호출 시 필요한 8개의 데이터를 전달합니다.
            ppt_file_path = create_cs_ppt_card(
                row, name_col, branch_col, q_text, a_text, dept_info, collab_dept, status_bg
            )
            
            with open(ppt_file_path, "rb") as f:
                st.download_button(
                    label="💾 생성된 PPT 파일 PC에 저장하기",
                    data=f,
                    file_name=f"CS리포트_{row[name_col]}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            st.success("✅ PPT 생성이 완료되었습니다!")
        except Exception as e:
            st.error(f"PPT 생성 중 오류가 발생했습니다: {e}")